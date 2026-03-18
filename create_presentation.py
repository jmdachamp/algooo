from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
DARK_NAVY = RGBColor(0x0B, 0x1D, 0x3A)
GOLD = RGBColor(0xC8, 0x96, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MED_BLUE = RGBColor(0x1A, 0x3A, 0x6B)
LIGHT_BLUE = RGBColor(0x2C, 0x5F, 0x9E)
ACCENT_GREEN = RGBColor(0x2E, 0x8B, 0x57)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, color, left=0, top=0, width=None, height=None):
    w = width or prs.slide_width
    h = height or prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_gold_bar(slide, left, top, width=Inches(1.5), height=Inches(0.06)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape


def add_bullet_slide(slide, items, left, top, width, height, font_size=18, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(10)
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"\u2022  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox


# ── SLIDE 1: Title Slide ──
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_NAVY)

# Gold accent bar at top
add_shape_bg(slide, GOLD, Inches(0), Inches(0), prs.slide_width, Inches(0.08))

add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.5),
            "LIFE INSURANCE", 54, GOLD, True, PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.1), Inches(11), Inches(1),
            "Your Trusted Broker for the Best Coverage & Rates", 28, WHITE, False, PP_ALIGN.CENTER)

add_gold_bar(slide, Inches(5.5), Inches(4.2), Inches(2.5), Inches(0.06))

add_textbox(slide, Inches(1), Inches(4.8), Inches(11), Inches(1),
            "We shop the market so you don't have to.", 22, LIGHT_GRAY, False, PP_ALIGN.CENTER)


# ── SLIDE 2: Why Life Insurance Matters ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_shape_bg(slide, GOLD, Inches(0), Inches(0), prs.slide_width, Inches(0.08))

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
            "WHY LIFE INSURANCE MATTERS", 36, GOLD, True, PP_ALIGN.LEFT)
add_gold_bar(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.05))

items = [
    "Provides financial security for your loved ones if the unexpected happens",
    "Replaces lost income to maintain your family\u2019s standard of living",
    "Covers outstanding debts \u2014 mortgages, loans, and credit cards",
    "Funds your children\u2019s education and future milestones",
    "Covers funeral and final expenses, easing the burden on family",
    "Offers peace of mind knowing your family is protected"
]
add_bullet_slide(slide, items, Inches(1), Inches(1.8), Inches(11), Inches(5), 22, WHITE)


# ── SLIDE 3: Who We Are ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_shape_bg(slide, GOLD, Inches(0), Inches(0), prs.slide_width, Inches(0.08))

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
            "WHO WE ARE", 36, GOLD, True, PP_ALIGN.LEFT)
add_gold_bar(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.05))

add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
            "We Are Independent Life Insurance Brokers", 26, WHITE, True, PP_ALIGN.LEFT)

items = [
    "We are NOT tied to any single insurance company",
    "We work for YOU \u2014 not the insurance carriers",
    "Access to dozens of top-rated insurance providers nationwide",
    "Licensed professionals with deep industry expertise",
    "Our mission: find you the best coverage at the best price"
]
add_bullet_slide(slide, items, Inches(1), Inches(2.8), Inches(11), Inches(4), 22, WHITE)


# ── SLIDE 4: The Broker Advantage ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_shape_bg(slide, GOLD, Inches(0), Inches(0), prs.slide_width, Inches(0.08))

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
            "THE BROKER ADVANTAGE", 36, GOLD, True, PP_ALIGN.LEFT)
add_gold_bar(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.05))

# Three columns
col_w = Inches(3.5)
col_h = Inches(4)
cols = [
    ("CHOICE", [
        "We compare policies from",
        "multiple carriers to find",
        "the perfect fit for your",
        "needs and budget."
    ]),
    ("SAVINGS", [
        "Competitive quoting means",
        "you get the best rates \u2014",
        "often saving 20-40% vs.",
        "going direct to one carrier."
    ]),
    ("ADVOCACY", [
        "We handle the paperwork,",
        "negotiate on your behalf,",
        "and support you through",
        "the entire claims process."
    ]),
]

for i, (title, lines) in enumerate(cols):
    x = Inches(0.8 + i * 4.2)
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), col_w, Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = MED_BLUE
    card.line.fill.background()

    add_textbox(slide, x + Inches(0.3), Inches(2.1), Inches(2.9), Inches(0.7),
                title, 24, GOLD, True, PP_ALIGN.CENTER)
    add_gold_bar(slide, x + Inches(1), Inches(2.8), Inches(1.5), Inches(0.04))

    body_text = "\n".join(lines)
    add_textbox(slide, x + Inches(0.3), Inches(3.1), Inches(2.9), Inches(3),
                body_text, 18, WHITE, False, PP_ALIGN.CENTER)


# ── SLIDE 5: Types of Life Insurance ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_shape_bg(slide, GOLD, Inches(0), Inches(0), prs.slide_width, Inches(0.08))

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
            "TYPES OF LIFE INSURANCE", 36, GOLD, True, PP_ALIGN.LEFT)
add_gold_bar(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.05))

types = [
    ("TERM LIFE", "Affordable coverage for a set period (10, 20, or 30 years). Ideal for young families and mortgage protection."),
    ("WHOLE LIFE", "Lifetime coverage with a guaranteed cash value component. Builds wealth over time with fixed premiums."),
    ("UNIVERSAL LIFE", "Flexible premiums and adjustable death benefits. Offers cash value growth tied to market performance."),
    ("INDEXED UNIVERSAL LIFE", "Cash value growth linked to a stock index with downside protection. Upside potential with a safety net."),
]

for i, (title, desc) in enumerate(types):
    y = Inches(1.7 + i * 1.35)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.5), Inches(1.15))
    card.fill.solid()
    card.fill.fore_color.rgb = MED_BLUE
    card.line.fill.background()

    add_textbox(slide, Inches(1.1), y + Inches(0.1), Inches(3), Inches(0.5),
                title, 20, GOLD, True, PP_ALIGN.LEFT)
    add_textbox(slide, Inches(1.1), y + Inches(0.55), Inches(10.8), Inches(0.6),
                desc, 16, WHITE, False, PP_ALIGN.LEFT)


# ── SLIDE 6: How We Get You the Best Deal ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_shape_bg(slide, GOLD, Inches(0), Inches(0), prs.slide_width, Inches(0.08))

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
            "HOW WE GET YOU THE BEST DEAL", 36, GOLD, True, PP_ALIGN.LEFT)
add_gold_bar(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.05))

steps = [
    ("1", "DISCOVER", "We learn about your needs, goals, health, and budget in a free consultation."),
    ("2", "SHOP", "We compare quotes from 30+ top-rated carriers to find the best options."),
    ("3", "RECOMMEND", "We present the top choices with clear comparisons \u2014 no jargon, no pressure."),
    ("4", "SECURE", "We handle the application, underwriting, and paperwork for a smooth process."),
    ("5", "SUPPORT", "We\u2019re here for you year after year \u2014 reviews, updates, and claims assistance."),
]

for i, (num, title, desc) in enumerate(steps):
    y = Inches(1.7 + i * 1.1)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y, Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = num
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = DARK_NAVY
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_textbox(slide, Inches(2), y + Inches(0.0), Inches(3), Inches(0.5),
                title, 22, GOLD, True, PP_ALIGN.LEFT)
    add_textbox(slide, Inches(2), y + Inches(0.4), Inches(10), Inches(0.5),
                desc, 17, WHITE, False, PP_ALIGN.LEFT)


# ── SLIDE 7: Why Choose Us ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_shape_bg(slide, GOLD, Inches(0), Inches(0), prs.slide_width, Inches(0.08))

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
            "WHY CLIENTS CHOOSE US", 36, GOLD, True, PP_ALIGN.LEFT)
add_gold_bar(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.05))

reasons = [
    "\u2713  Independent & unbiased \u2014 we work for you, not insurance companies",
    "\u2713  Access to 30+ top-rated carriers for the most competitive rates",
    "\u2713  Personalized guidance tailored to your unique situation",
    "\u2713  No cost to you \u2014 our services are free (carriers pay our commission)",
    "\u2713  Expert knowledge of underwriting to match you with the right carrier",
    "\u2713  Ongoing support \u2014 annual reviews and claims advocacy",
    "\u2713  Fast, hassle-free process from quote to coverage"
]
add_bullet_slide(slide, reasons, Inches(1), Inches(1.8), Inches(11), Inches(5), 22, WHITE)


# ── SLIDE 8: Client Testimonials ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_shape_bg(slide, GOLD, Inches(0), Inches(0), prs.slide_width, Inches(0.08))

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
            "WHAT OUR CLIENTS SAY", 36, GOLD, True, PP_ALIGN.LEFT)
add_gold_bar(slide, Inches(0.8), Inches(1.3), Inches(2), Inches(0.05))

testimonials = [
    ("\u201cThey saved us over $1,200 a year by shopping multiple carriers. We never would have found that rate on our own.\u201d", "\u2014 The Johnson Family"),
    ("\u201cThe process was so easy. They handled everything and I had coverage within a week.\u201d", "\u2014 Michael R."),
    ("\u201cAs a small business owner, I needed the right policy to protect my family and my business. They found the perfect solution.\u201d", "\u2014 Sarah T., Business Owner"),
]

for i, (quote, author) in enumerate(testimonials):
    y = Inches(1.7 + i * 1.8)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(10), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = MED_BLUE
    card.line.fill.background()

    add_textbox(slide, Inches(2), y + Inches(0.15), Inches(9), Inches(0.8),
                quote, 18, WHITE, False, PP_ALIGN.LEFT)
    add_textbox(slide, Inches(2), y + Inches(0.95), Inches(9), Inches(0.4),
                author, 16, GOLD, True, PP_ALIGN.RIGHT)


# ── SLIDE 9: Get Started / Contact ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)
add_shape_bg(slide, GOLD, Inches(0), Inches(0), prs.slide_width, Inches(0.08))

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
            "GET YOUR FREE QUOTE TODAY", 44, GOLD, True, PP_ALIGN.CENTER)

add_gold_bar(slide, Inches(5.5), Inches(2.8), Inches(2.5), Inches(0.06))

add_textbox(slide, Inches(1), Inches(3.3), Inches(11), Inches(1),
            "No obligation. No pressure. Just honest advice and the best rates.", 24, WHITE, False, PP_ALIGN.CENTER)

# Contact info card
card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(4.3), Inches(6), Inches(2.2))
card.fill.solid()
card.fill.fore_color.rgb = MED_BLUE
card.line.fill.background()

contact_lines = [
    "Contact Us",
    "",
    "Phone: (XXX) XXX-XXXX",
    "Email: info@yourcompany.com",
    "Web: www.yourcompany.com",
]

txBox = slide.shapes.add_textbox(Inches(3.8), Inches(4.5), Inches(5.4), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
for i, line in enumerate(contact_lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = line
    if i == 0:
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = GOLD
    else:
        run.font.size = Pt(18)
        run.font.color.rgb = WHITE
    run.font.name = "Calibri"

add_textbox(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
            "Let us do the shopping \u2014 so you get the best protection at the best price.", 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)


# Save
prs.save("Life_Insurance_Presentation.pptx")
print("Presentation saved: Life_Insurance_Presentation.pptx")
